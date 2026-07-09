"""
文档解析模块
支持 PDF、DOCX、PPTX、TXT、MD 格式的文档解析与分块
"""

import os
import logging
from typing import List

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import settings

logger = logging.getLogger(__name__)


def _parse_pdf(file_path: str) -> List[Document]:
    """
    使用 pymupdf (fitz) 解析 PDF 文件

    Args:
        file_path: PDF 文件路径

    Returns:
        解析后的 Document 列表，每页一个 Document
    """
    import fitz

    documents = []
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                documents.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "page": page_num + 1,
                            "file_type": "pdf",
                        },
                    )
                )
        doc.close()
    except Exception as e:
        logger.error(f"解析 PDF 文件失败 {file_path}: {e}")
        raise
    return documents


def _parse_docx(file_path: str) -> List[Document]:
    """
    使用 python-docx 解析 DOCX 文件

    Args:
        file_path: DOCX 文件路径

    Returns:
        解析后的 Document 列表
    """
    from docx import Document as DocxDocument

    documents = []
    try:
        doc = DocxDocument(file_path)
        full_text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text_parts.append(para.text)
        full_text = "\n".join(full_text_parts)
        if full_text.strip():
            documents.append(
                Document(
                    page_content=full_text,
                    metadata={
                        "source": file_path,
                        "page": 1,
                        "file_type": "docx",
                    },
                )
            )
    except Exception as e:
        logger.error(f"解析 DOCX 文件失败 {file_path}: {e}")
        raise
    return documents


def _parse_pptx(file_path: str) -> List[Document]:
    """
    使用 python-pptx 解析 PPTX 文件

    Args:
        file_path: PPTX 文件路径

    Returns:
        解析后的 Document 列表，每页幻灯片一个 Document
    """
    from pptx import Presentation

    documents = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slide_text = "\n".join(texts)
                documents.append(
                    Document(
                        page_content=slide_text,
                        metadata={
                            "source": file_path,
                            "page": slide_num,
                            "file_type": "pptx",
                        },
                    )
                )
    except Exception as e:
        logger.error(f"解析 PPTX 文件失败 {file_path}: {e}")
        raise
    return documents


def _parse_text(file_path: str) -> List[Document]:
    """
    解析纯文本文件（TXT/MD）

    Args:
        file_path: 文本文件路径

    Returns:
        解析后的 Document 列表
    """
    documents = []
    try:
        # 尝试多种编码，优先 UTF-8
        content = None
        for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        if content is None:
            raise ValueError(f"无法解码文件 {file_path}，尝试了多种编码均失败")

        if content.strip():
            ext = os.path.splitext(file_path)[1].lower()
            file_type = "markdown" if ext in (".md", ".markdown") else "txt"
            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "page": 1,
                        "file_type": file_type,
                    },
                )
            )
    except Exception as e:
        logger.error(f"解析文本文件失败 {file_path}: {e}")
        raise
    return documents


# 文件扩展名到解析函数的映射
_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt": _parse_text,
    ".md": _parse_text,
    ".markdown": _parse_text,
}


def parse_document(file_path: str) -> List[Document]:
    """
    解析单个文档文件，根据文件类型选择对应的解析器

    Args:
        file_path: 文件路径

    Returns:
        解析后的 Document 列表，包含元数据（source, page, file_type）

    Raises:
        ValueError: 不支持的文件类型
        FileNotFoundError: 文件不存在
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(
            f"不支持的文件类型: {ext}，"
            f"支持的格式: {', '.join(_PARSERS.keys())}"
        )

    logger.info(f"开始解析文档: {file_path} (类型: {ext})")
    documents = parser(file_path)
    logger.info(f"文档解析完成: {file_path}，共 {len(documents)} 页/段")
    return documents


def chunk_documents(documents: List[Document]) -> List[Document]:
    """
    使用 RecursiveCharacterTextSplitter 对文档进行分块

    Args:
        documents: 原始 Document 列表

    Returns:
        分块后的 Document 列表，保留原始元数据
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", "。", "！", "？", "；", "，", " ", ""],
    )

    chunks = text_splitter.split_documents(documents)
    logger.info(f"文档分块完成：{len(documents)} 个文档 -> {len(chunks)} 个分块")
    return chunks


def parse_and_chunk(file_path: str) -> List[Document]:
    """
    解析文档并自动分块，一站式处理

    Args:
        file_path: 文件路径

    Returns:
        分块后的 Document 列表
    """
    documents = parse_document(file_path)
    return chunk_documents(documents)
